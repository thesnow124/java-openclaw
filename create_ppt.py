#!/usr/bin/env python3
from pptx import Presentation

# 创建演示文稿
pr = Presentation()

# 添加标题幻灯片
slide = pr.slides.add_slide(pr.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "test789"
subtitle.text = "创建于 2026-02-07"

# 保存文件
pr.save('test789.pptx')
print('PPT 文件已创建: test789.pptx')
